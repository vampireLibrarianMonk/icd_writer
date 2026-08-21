"""Tests for corpus document generation — format and presentation validation.

Verifies that generated DOCX, PPTX, XLSX, HTML, and image files:
1. Are valid files that open without errors
2. Contain real aerospace content (not placeholder/lorem ipsum)
3. Have professional formatting (fonts, headers, structure)
4. Look legitimate as engineering source documents

These tests run against the generated corpus in test_corpus/.
Run the generation script first: python test_corpus/scripts/generate_corpus.py
"""

from pathlib import Path

import pytest

CORPUS_DIR = Path(__file__).parent.parent.parent.parent / "test_corpus"
HSI_DIR = CORPUS_DIR / "hsi_sys_015g"
TSAFE_DIR = CORPUS_DIR / "tsafe"


def _corpus_exists():
    return HSI_DIR.exists() and (HSI_DIR / "docx").exists()


pytestmark = pytest.mark.skipif(
    not _corpus_exists(),
    reason="Test corpus not generated (run: python test_corpus/scripts/generate_corpus.py)",
)


# ═══════════════════════════════════════════════════════════════════════
# DOCX VALIDATION
# ═══════════════════════════════════════════════════════════════════════


class TestDocxFormat:
    """Validate DOCX files are properly formatted and contain real content."""

    @pytest.fixture(params=[
        "hsi_sys_015g/docx/HSI_Mech_Requirements_v2.docx",
        "hsi_sys_015g/docx/Thermal_Test_Report.docx",
        "hsi_sys_015g/docx/Integration_Procedure_v1.docx",
        "tsafe/docx/TSAFE_Message_Spec.docx",
    ])
    def docx_path(self, request):
        return CORPUS_DIR / request.param

    def test_docx_opens_without_error(self, docx_path):
        """File opens as a valid DOCX document."""
        from docx import Document
        doc = Document(str(docx_path))
        assert len(doc.paragraphs) > 0

    def test_docx_has_title(self, docx_path):
        """Document has a title heading (not empty first page)."""
        from docx import Document
        doc = Document(str(docx_path))
        # Find first heading
        headings = [p for p in doc.paragraphs if p.style.name.startswith("Heading")]
        assert len(headings) >= 1, "No headings found in document"

    def test_docx_has_multiple_sections(self, docx_path):
        """Document has multiple section headings (structured, not flat text)."""
        from docx import Document
        doc = Document(str(docx_path))
        headings = [p for p in doc.paragraphs if p.style.name.startswith("Heading")]
        assert len(headings) >= 3, f"Only {len(headings)} headings — too few for a spec doc"

    def test_docx_body_has_substance(self, docx_path):
        """Body paragraphs contain substantial text (not stubs)."""
        from docx import Document
        doc = Document(str(docx_path))
        body_paras = [
            p for p in doc.paragraphs
            if p.style.name == "Normal" and len(p.text.strip()) > 20
        ]
        assert len(body_paras) >= 5, f"Only {len(body_paras)} substantial paragraphs"

    def test_docx_contains_aerospace_terms(self, docx_path):
        """Content contains real aerospace/engineering terminology."""
        from docx import Document
        doc = Document(str(docx_path))
        full_text = " ".join(p.text for p in doc.paragraphs).lower()

        aerospace_terms = [
            "interface", "specification", "requirement", "system",
            "design", "test", "configuration", "performance",
            "data", "format", "record", "protocol", "conflict",
            "track", "radar", "aircraft", "advisory", "message",
            "thermal", "mechanical", "electrical", "power",
        ]
        matches = [t for t in aerospace_terms if t in full_text]
        assert len(matches) >= 3, (
            f"Only found terms {matches} — content doesn't look like engineering doc"
        )

    def test_docx_no_placeholder_text(self, docx_path):
        """No lorem ipsum or placeholder markers in the text."""
        from docx import Document
        doc = Document(str(docx_path))
        full_text = " ".join(p.text for p in doc.paragraphs).lower()

        placeholders = ["lorem ipsum", "todo", "placeholder", "insert text here", "tbd content"]
        for ph in placeholders:
            assert ph not in full_text, f"Placeholder '{ph}' found in {docx_path.name}"

    def test_docx_has_professional_font(self, docx_path):
        """Default font is a professional serif or sans-serif (not Comic Sans)."""
        from docx import Document
        doc = Document(str(docx_path))
        style = doc.styles["Normal"]
        font_name = style.font.name
        professional_fonts = [
            "Times New Roman", "Arial", "Calibri", "Cambria",
            "Helvetica", "Georgia", None,  # None = inherits default
        ]
        assert font_name in professional_fonts, f"Unprofessional font: {font_name}"

    def test_docx_file_size_reasonable(self, docx_path):
        """File size is reasonable (not empty, not bloated)."""
        size = docx_path.stat().st_size
        assert size > 5_000, f"Too small ({size} bytes) — likely empty"
        assert size < 5_000_000, f"Too large ({size} bytes) — bloated"


# ═══════════════════════════════════════════════════════════════════════
# PPTX VALIDATION
# ═══════════════════════════════════════════════════════════════════════


class TestPptxFormat:
    """Validate PPTX files look like real engineering presentations."""

    @pytest.fixture(params=[
        "hsi_sys_015g/pptx/Thermal_Analysis_CDR.pptx",
        "tsafe/pptx/TSAFE_Architecture_Overview.pptx",
    ])
    def pptx_path(self, request):
        return CORPUS_DIR / request.param

    def test_pptx_opens_without_error(self, pptx_path):
        """File opens as a valid PPTX presentation."""
        from pptx import Presentation
        prs = Presentation(str(pptx_path))
        assert len(prs.slides) > 0

    def test_pptx_has_title_slide(self, pptx_path):
        """First slide has a title (not blank)."""
        from pptx import Presentation
        prs = Presentation(str(pptx_path))
        first_slide = prs.slides[0]
        title_shape = first_slide.shapes.title
        assert title_shape is not None
        assert len(title_shape.text.strip()) > 5, "Title slide text too short"

    def test_pptx_has_multiple_slides(self, pptx_path):
        """Presentation has at least 4 slides (title + content + summary)."""
        from pptx import Presentation
        prs = Presentation(str(pptx_path))
        assert len(prs.slides) >= 4, f"Only {len(prs.slides)} slides — too few for a review"

    def test_pptx_slides_have_content(self, pptx_path):
        """Content slides have text (not blank slides)."""
        from pptx import Presentation
        prs = Presentation(str(pptx_path))

        empty_slides = 0
        for i, slide in enumerate(prs.slides):
            if i == 0:
                continue  # Skip title slide
            text = " ".join(
                shape.text for shape in slide.shapes if shape.has_text_frame
            )
            if len(text.strip()) < 10:
                empty_slides += 1

        assert empty_slides <= 1, f"{empty_slides} slides are essentially empty"

    def test_pptx_contains_technical_content(self, pptx_path):
        """Slides contain technical/engineering language."""
        from pptx import Presentation
        prs = Presentation(str(pptx_path))
        all_text = " ".join(
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
        ).lower()

        technical_terms = [
            "thermal", "power", "interface", "design", "system",
            "performance", "test", "requirement", "data", "architecture",
        ]
        matches = [t for t in technical_terms if t in all_text]
        assert len(matches) >= 3, f"Only found terms {matches}"

    def test_pptx_file_size_reasonable(self, pptx_path):
        """File size is reasonable for a presentation."""
        size = pptx_path.stat().st_size
        assert size > 10_000, f"Too small ({size} bytes)"
        assert size < 10_000_000, f"Too large ({size} bytes)"


# ═══════════════════════════════════════════════════════════════════════
# XLSX VALIDATION
# ═══════════════════════════════════════════════════════════════════════


class TestXlsxFormat:
    """Validate XLSX files have proper structure and engineering data."""

    @pytest.fixture(params=[
        "hsi_sys_015g/xlsx/Thermostat_Parameters.xlsx",
        "hsi_sys_015g/xlsx/HSI_Mass_Properties.xlsx",
        "hsi_sys_015g/xlsx/Power_Budget_v2.xlsx",
        "tsafe/xlsx/TSAFE_Test_Results.xlsx",
    ])
    def xlsx_path(self, request):
        return CORPUS_DIR / request.param

    def test_xlsx_opens_without_error(self, xlsx_path):
        """File opens as a valid XLSX workbook."""
        from openpyxl import load_workbook
        wb = load_workbook(str(xlsx_path))
        assert len(wb.sheetnames) >= 1

    def test_xlsx_has_header_row(self, xlsx_path):
        """Worksheet has a formatted header row (not raw data starting at row 1)."""
        from openpyxl import load_workbook
        wb = load_workbook(str(xlsx_path))
        ws = wb.active

        # Find the header row (row with bold or filled cells)
        header_row = None
        for row in range(1, 10):
            cell = ws.cell(row=row, column=1)
            if cell.font and cell.font.bold:
                header_row = row
                break
            if cell.fill and cell.fill.fgColor and cell.fill.fgColor.rgb != "00000000":
                header_row = row
                break

        assert header_row is not None, "No formatted header row found"

    def test_xlsx_has_data_rows(self, xlsx_path):
        """Worksheet has at least 3 data rows below headers."""
        from openpyxl import load_workbook
        wb = load_workbook(str(xlsx_path))
        ws = wb.active

        # Count non-empty rows
        data_rows = 0
        for row in ws.iter_rows(min_row=5, max_row=50, max_col=1):
            if row[0].value:
                data_rows += 1

        assert data_rows >= 3, f"Only {data_rows} data rows — too sparse"

    def test_xlsx_has_numeric_values(self, xlsx_path):
        """At least some cells contain numeric values (not all text)."""
        from openpyxl import load_workbook
        wb = load_workbook(str(xlsx_path))
        ws = wb.active

        numeric_cells = 0
        for row in ws.iter_rows(min_row=5, max_row=20, min_col=2, max_col=7):
            for cell in row:
                if cell.value is not None:
                    try:
                        float(str(cell.value).replace(",", ""))
                        numeric_cells += 1
                    except (ValueError, TypeError):
                        pass

        assert numeric_cells >= 3, f"Only {numeric_cells} numeric cells — should have parameters"

    def test_xlsx_has_title(self, xlsx_path):
        """Workbook has a title in the first row (merged or large font)."""
        from openpyxl import load_workbook
        wb = load_workbook(str(xlsx_path))
        ws = wb.active
        title_cell = ws.cell(row=1, column=1)
        assert title_cell.value is not None, "No title in cell A1"
        assert len(str(title_cell.value)) > 5, "Title too short"

    def test_xlsx_has_borders(self, xlsx_path):
        """Data cells have borders (professional formatting)."""
        from openpyxl import load_workbook
        wb = load_workbook(str(xlsx_path))
        ws = wb.active

        bordered_cells = 0
        for row in ws.iter_rows(min_row=4, max_row=10, min_col=1, max_col=5):
            for cell in row:
                if cell.border and (cell.border.left.style or cell.border.top.style):
                    bordered_cells += 1

        assert bordered_cells >= 5, f"Only {bordered_cells} bordered cells — missing formatting"

    def test_xlsx_file_size_reasonable(self, xlsx_path):
        """File size is reasonable."""
        size = xlsx_path.stat().st_size
        assert size > 3_000, f"Too small ({size} bytes)"
        assert size < 2_000_000, f"Too large ({size} bytes)"


# ═══════════════════════════════════════════════════════════════════════
# HTML VALIDATION
# ═══════════════════════════════════════════════════════════════════════


class TestHtmlFormat:
    """Validate HTML files look like Confluence wiki pages."""

    @pytest.fixture(params=[
        "hsi_sys_015g/html/HSI_Power_ICD_Wiki.html",
    ])
    def html_path(self, request):
        return CORPUS_DIR / request.param

    def test_html_is_valid(self, html_path):
        """File is valid HTML with proper structure."""
        content = html_path.read_text(encoding="utf-8")
        assert "<!DOCTYPE html>" in content
        assert "<html>" in content or "<html " in content
        assert "</html>" in content
        assert "<body>" in content

    def test_html_has_title(self, html_path):
        """Page has a meaningful title."""
        content = html_path.read_text(encoding="utf-8")
        assert "<title>" in content
        # Extract title
        start = content.index("<title>") + 7
        end = content.index("</title>")
        title = content[start:end]
        assert len(title) > 5, f"Title too short: '{title}'"

    def test_html_has_headings(self, html_path):
        """Page has section headings (h2 or h3 tags)."""
        content = html_path.read_text(encoding="utf-8")
        h2_count = content.count("<h2>") + content.count("<h2 ")
        h3_count = content.count("<h3>") + content.count("<h3 ")
        assert h2_count + h3_count >= 2, "Too few section headings"

    def test_html_has_body_content(self, html_path):
        """Page has substantial paragraph content."""
        content = html_path.read_text(encoding="utf-8")
        p_count = content.count("<p>") + content.count("<p ")
        assert p_count >= 4, f"Only {p_count} paragraphs"

    def test_html_has_confluence_styling(self, html_path):
        """Page has CSS that looks like a wiki/Confluence page."""
        content = html_path.read_text(encoding="utf-8")
        assert "<style>" in content
        # Should have sans-serif body font (wiki style)
        assert "sans-serif" in content or "Segoe UI" in content

    def test_html_has_metadata(self, html_path):
        """Page shows author, date, or space information."""
        content = html_path.read_text(encoding="utf-8")
        # Should have some metadata indicators
        has_date = "modified" in content.lower() or "2025" in content or "2024" in content
        has_author = "author" in content.lower() or "engineering" in content.lower()
        assert has_date or has_author, "No metadata (date/author) visible"

    def test_html_contains_technical_content(self, html_path):
        """Content has engineering terminology, not filler."""
        content = html_path.read_text(encoding="utf-8").lower()
        terms = ["power", "interface", "voltage", "current", "mode", "requirement"]
        matches = [t for t in terms if t in content]
        assert len(matches) >= 3, f"Only found: {matches}"


# ═══════════════════════════════════════════════════════════════════════
# IMAGE VALIDATION
# ═══════════════════════════════════════════════════════════════════════


class TestImageFormat:
    """Validate images are non-trivial and properly formatted."""

    @pytest.fixture(params=[
        "hsi_sys_015g/images/Spectrometer_Mount_Drawing.png",
        "hsi_sys_015g/images/Thermal_Table_Extract.png",
        "hsi_sys_015g/images/System_Overview_Diagram.png",
        "tsafe/images/TSAFE_Architecture_Diagram.png",
        "tsafe/images/Radar_Interface_Diagram.png",
    ])
    def png_path(self, request):
        return CORPUS_DIR / request.param

    def test_png_valid_format(self, png_path):
        """File has valid PNG magic bytes."""
        data = png_path.read_bytes()
        assert data[:8] == b"\x89PNG\r\n\x1a\n", "Not a valid PNG file"

    def test_png_non_trivial_size(self, png_path):
        """Image is large enough to contain actual content (not a 1x1 pixel)."""
        size = png_path.stat().st_size
        assert size > 10_000, f"Image too small ({size} bytes) — likely blank"

    def test_png_reasonable_dimensions(self, png_path):
        """Image has reasonable dimensions (not microscopic or absurdly large)."""
        from PIL import Image
        img = Image.open(str(png_path))
        w, h = img.size
        assert w >= 200, f"Width {w}px too small"
        assert h >= 150, f"Height {h}px too small"
        assert w <= 5000, f"Width {w}px too large"
        assert h <= 5000, f"Height {h}px too large"

    def test_png_not_solid_color(self, png_path):
        """Image is not a solid single color (has actual content)."""
        from PIL import Image
        import statistics

        img = Image.open(str(png_path)).convert("L")  # Grayscale
        pixels = list(img.getdata())
        # Standard deviation of pixel values — solid color = 0
        stdev = statistics.stdev(pixels[:10000])  # Sample first 10K pixels
        assert stdev > 5, f"Image appears to be solid color (stdev={stdev:.1f})"


class TestTiffFormat:
    """Validate TIFF files simulate scanned drawings."""

    @pytest.fixture(params=[
        "hsi_sys_015g/images/Connector_Pinout_J1.tiff",
    ])
    def tiff_path(self, request):
        return CORPUS_DIR / request.param

    def test_tiff_valid_format(self, tiff_path):
        """File opens as a valid TIFF image."""
        from PIL import Image
        img = Image.open(str(tiff_path))
        assert img.format == "TIFF"

    def test_tiff_is_grayscale(self, tiff_path):
        """Scanned drawing is grayscale (simulating a real scan)."""
        from PIL import Image
        img = Image.open(str(tiff_path))
        assert img.mode in ("L", "1", "LA"), f"Expected grayscale, got mode={img.mode}"

    def test_tiff_reasonable_dimensions(self, tiff_path):
        """TIFF has page-like dimensions (letter/A4 at scan DPI)."""
        from PIL import Image
        img = Image.open(str(tiff_path))
        w, h = img.size
        # At 200 DPI, letter size = ~1700x2200
        assert w >= 800, f"Width {w}px too narrow for a scanned page"
        assert h >= 1000, f"Height {h}px too short for a scanned page"

    def test_tiff_has_content(self, tiff_path):
        """TIFF has actual content (not blank white page)."""
        from PIL import Image
        import statistics

        img = Image.open(str(tiff_path)).convert("L")
        pixels = list(img.getdata())
        # Sample broadly across the image
        step = max(1, len(pixels) // 50000)
        sample = pixels[::step][:50000]
        stdev = statistics.stdev(sample)
        assert stdev > 5, f"TIFF appears blank (stdev={stdev:.1f})"


# ═══════════════════════════════════════════════════════════════════════
# CROSS-CUTTING: CORPUS COMPLETENESS
# ═══════════════════════════════════════════════════════════════════════


class TestCorpusCompleteness:
    """Validate the overall corpus has expected file counts and diversity."""

    def test_hsi_has_all_format_types(self):
        """HSI corpus has DOCX, XLSX, PPTX, HTML, and images."""
        assert (HSI_DIR / "docx").exists()
        assert (HSI_DIR / "xlsx").exists()
        assert (HSI_DIR / "pptx").exists()
        assert (HSI_DIR / "html").exists()
        assert (HSI_DIR / "images").exists()

    def test_hsi_docx_count(self):
        """HSI has at least 3 DOCX files."""
        docx_files = list((HSI_DIR / "docx").glob("*.docx"))
        assert len(docx_files) >= 3, f"Only {len(docx_files)} DOCX files"

    def test_hsi_xlsx_count(self):
        """HSI has at least 3 XLSX files."""
        xlsx_files = list((HSI_DIR / "xlsx").glob("*.xlsx"))
        assert len(xlsx_files) >= 3, f"Only {len(xlsx_files)} XLSX files"

    def test_hsi_image_count(self):
        """HSI has at least 3 image files (PNG + TIFF)."""
        images = list((HSI_DIR / "images").glob("*.png"))
        images += list((HSI_DIR / "images").glob("*.tiff"))
        images += list((HSI_DIR / "images").glob("*.jpg"))
        assert len(images) >= 3, f"Only {len(images)} image files"

    def test_tsafe_has_multiple_formats(self):
        """TSAFE corpus has at least 3 different format types."""
        formats_present = 0
        if list((TSAFE_DIR / "docx").glob("*.docx")):
            formats_present += 1
        if list((TSAFE_DIR / "xlsx").glob("*.xlsx")):
            formats_present += 1
        if list((TSAFE_DIR / "pptx").glob("*.pptx")):
            formats_present += 1
        if (TSAFE_DIR / "images").exists() and list((TSAFE_DIR / "images").glob("*")):
            formats_present += 1
        assert formats_present >= 3, f"Only {formats_present} format types present"

    def test_manifest_exists(self):
        """Corpus manifest JSON exists."""
        manifest = CORPUS_DIR / "corpus_manifest.json"
        assert manifest.exists(), "corpus_manifest.json not found"

    def test_manifest_valid_json(self):
        """Manifest is valid JSON with expected fields."""
        import json
        manifest = CORPUS_DIR / "corpus_manifest.json"
        data = json.loads(manifest.read_text(encoding="utf-8"))
        assert "generated_at" in data
        assert "source_icds" in data
        assert len(data["source_icds"]) >= 5
