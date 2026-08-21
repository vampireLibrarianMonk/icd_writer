"""Generate the test corpus from existing ICD Document IRs.

Creates realistic DOCX, PPTX, XLSX, HTML, and image files derived from
the actual content in our NASA ICD PDFs. Each generated file represents
an upstream source document that would feed into the ICD.

Usage:
    python test_corpus/scripts/generate_corpus.py

Output goes to test_corpus/hsi_sys_015g/, test_corpus/idss_idd/, etc.
"""

from __future__ import annotations

import json
import shutil
from datetime import datetime
from pathlib import Path

import fitz

# ─── Paths ─────────────────────────────────────────────────────────────

ROOT = Path(__file__).parent.parent.parent
ICDS_DIR = ROOT / "icds" / "digital"
CORPUS_DIR = Path(__file__).parent.parent
MANIFEST_PATH = CORPUS_DIR / "corpus_manifest.json"


# ─── Text Extraction Helpers ───────────────────────────────────────────


def extract_page_text(pdf_path: Path, page_num: int) -> str:
    """Extract full text from a single page (1-indexed)."""
    doc = fitz.open(str(pdf_path))
    text = doc[page_num - 1].get_text("text")
    doc.close()
    return text


def extract_pages_text(pdf_path: Path, page_range: range) -> str:
    """Extract text from a range of pages (1-indexed)."""
    doc = fitz.open(str(pdf_path))
    text = ""
    for i in page_range:
        text += doc[i - 1].get_text("text") + "\n"
    doc.close()
    return text


def extract_page_image(pdf_path: Path, page_num: int, dpi: int = 150) -> bytes:
    """Render a page as PNG bytes."""
    doc = fitz.open(str(pdf_path))
    page = doc[page_num - 1]
    pix = page.get_pixmap(dpi=dpi)
    png_bytes = pix.tobytes("png")
    doc.close()
    return png_bytes


def extract_page_region_image(
    pdf_path: Path, page_num: int, rect: tuple[float, float, float, float], dpi: int = 200
) -> bytes:
    """Render a region of a page as PNG. rect = (x0, y0, x1, y1) in points."""
    doc = fitz.open(str(pdf_path))
    page = doc[page_num - 1]
    clip = fitz.Rect(*rect)
    pix = page.get_pixmap(dpi=dpi, clip=clip)
    png_bytes = pix.tobytes("png")
    doc.close()
    return png_bytes


# ─── DOCX Generation ──────────────────────────────────────────────────


def generate_docx(
    title: str,
    sections: list[tuple[str, str]],
    output_path: Path,
    author: str = "HESSI Systems Engineering",
    doc_number: str = "",
    revision: str = "v1",
    date: str = "2024-06-01",
):
    """Generate a professional-looking Word document.

    Args:
        title: Document title
        sections: List of (heading, body_text) tuples
        output_path: Where to save the .docx
        author: Document author/organization
        doc_number: Optional document number
        revision: Version string
        date: Document date
    """
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # Set default font
    style = doc.styles["Normal"]
    font = style.font
    font.name = "Times New Roman"
    font.size = Pt(11)

    # Title page content
    doc.add_paragraph()  # spacing
    title_para = doc.add_heading(title, level=0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    if doc_number:
        num_para = doc.add_paragraph(doc_number)
        num_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    meta_para = doc.add_paragraph(f"Revision: {revision}    Date: {date}")
    meta_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    author_para = doc.add_paragraph(f"Prepared by: {author}")
    author_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_page_break()

    # Table of Contents placeholder
    doc.add_heading("Table of Contents", level=1)
    for i, (heading, _) in enumerate(sections, 1):
        doc.add_paragraph(f"{i}. {heading}", style="List Number")
    doc.add_page_break()

    # Sections
    for heading, body in sections:
        doc.add_heading(heading, level=1)
        # Split body into paragraphs
        for para_text in body.split("\n"):
            para_text = para_text.strip()
            if para_text:
                doc.add_paragraph(para_text)

    # Save
    output_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(output_path))


# ─── PPTX Generation ──────────────────────────────────────────────────


def generate_pptx(
    title: str,
    subtitle: str,
    slides: list[tuple[str, str]],
    output_path: Path,
    date: str = "2024-06-01",
    author: str = "HESSI Systems Engineering",
):
    """Generate a professional presentation.

    Args:
        title: Presentation title (title slide)
        subtitle: Subtitle for title slide
        slides: List of (slide_title, bullet_content) tuples
        output_path: Where to save the .pptx
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = f"{subtitle}\n{author}\n{date}"

    # Content slides
    bullet_layout = prs.slide_layouts[1]  # Title and Content
    for slide_title, content in slides:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = slide_title

        tf = slide.placeholders[1].text_frame
        tf.clear()

        # Split content into bullet points
        lines = [l.strip() for l in content.split("\n") if l.strip()]
        for i, line in enumerate(lines):
            if i == 0:
                tf.paragraphs[0].text = line
            else:
                p = tf.add_paragraph()
                p.text = line
                p.level = 1 if line.startswith("-") or line.startswith("•") else 0

    # Summary/Questions slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Summary & Next Steps"
    tf = slide.placeholders[1].text_frame
    tf.text = "Questions?"

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


# ─── XLSX Generation ──────────────────────────────────────────────────


def generate_xlsx(
    title: str,
    headers: list[str],
    rows: list[list[str]],
    output_path: Path,
    sheet_name: str = "Parameters",
    doc_ref: str = "",
    date: str = "2024-06-01",
):
    """Generate a professional spreadsheet with headers and data.

    Args:
        title: Workbook title (in header row)
        headers: Column headers
        rows: Data rows
        output_path: Where to save
    """
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, Border, Side, PatternFill

    wb = Workbook()
    ws = wb.active
    ws.title = sheet_name

    # Title row
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=len(headers))
    title_cell = ws.cell(row=1, column=1, value=title)
    title_cell.font = Font(name="Arial", size=14, bold=True)
    title_cell.alignment = Alignment(horizontal="center")

    # Metadata row
    ws.merge_cells(start_row=2, start_column=1, end_row=2, end_column=len(headers))
    meta_cell = ws.cell(row=2, column=1, value=f"Reference: {doc_ref}  |  Date: {date}")
    meta_cell.font = Font(name="Arial", size=9, italic=True)

    # Header row
    header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
    header_font = Font(name="Arial", size=10, bold=True, color="FFFFFF")
    thin_border = Border(
        left=Side(style="thin"),
        right=Side(style="thin"),
        top=Side(style="thin"),
        bottom=Side(style="thin"),
    )

    for col, header in enumerate(headers, 1):
        cell = ws.cell(row=4, column=col, value=header)
        cell.font = header_font
        cell.fill = header_fill
        cell.border = thin_border
        cell.alignment = Alignment(horizontal="center", wrap_text=True)

    # Data rows
    data_font = Font(name="Arial", size=10)
    for row_idx, row_data in enumerate(rows, 5):
        for col_idx, value in enumerate(row_data, 1):
            cell = ws.cell(row=row_idx, column=col_idx, value=value)
            cell.font = data_font
            cell.border = thin_border
            cell.alignment = Alignment(wrap_text=True)

    # Auto-width columns
    for col in range(1, len(headers) + 1):
        max_len = max(
            len(str(ws.cell(row=r, column=col).value or ""))
            for r in range(4, len(rows) + 5)
        )
        ws.column_dimensions[chr(64 + col)].width = min(max_len + 4, 40)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(output_path))


# ─── HTML (Confluence page) Generation ─────────────────────────────────


def generate_html(
    title: str,
    sections: list[tuple[str, str]],
    output_path: Path,
    space: str = "HESSI Engineering",
    last_modified: str = "2025-05-30",
    author: str = "Systems Engineering",
):
    """Generate a Confluence-style wiki page as HTML.

    Args:
        title: Page title
        sections: List of (heading, body) tuples
        output_path: Where to save the .html
    """
    sections_html = ""
    for heading, body in sections:
        paras = "\n".join(
            f"        <p>{p.strip()}</p>" for p in body.split("\n") if p.strip()
        )
        sections_html += f"""
    <h2>{heading}</h2>
{paras}
"""

    html = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>{title}</title>
    <style>
        body {{
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
            max-width: 900px;
            margin: 40px auto;
            padding: 0 20px;
            line-height: 1.6;
            color: #172B4D;
        }}
        .page-header {{
            border-bottom: 2px solid #0052CC;
            padding-bottom: 12px;
            margin-bottom: 24px;
        }}
        .page-header h1 {{
            margin: 0;
            color: #172B4D;
        }}
        .page-meta {{
            font-size: 12px;
            color: #6B778C;
            margin-top: 8px;
        }}
        h2 {{
            color: #172B4D;
            border-bottom: 1px solid #DFE1E6;
            padding-bottom: 8px;
            margin-top: 32px;
        }}
        p {{
            margin: 12px 0;
        }}
        .status-badge {{
            display: inline-block;
            padding: 2px 8px;
            border-radius: 3px;
            font-size: 11px;
            font-weight: bold;
            background: #00875A;
            color: white;
        }}
    </style>
</head>
<body>
    <div class="page-header">
        <h1>{title}</h1>
        <div class="page-meta">
            <span>Space: {space}</span> &middot;
            <span>Last modified: {last_modified}</span> &middot;
            <span>Author: {author}</span> &middot;
            <span class="status-badge">CURRENT</span>
        </div>
    </div>
{sections_html}
    <hr>
    <p style="font-size: 11px; color: #6B778C;">
        This page is maintained as a living specification. Changes are tracked
        in the page history. For the formal baseline, see the released ICD PDF.
    </p>
</body>
</html>"""

    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(html, encoding="utf-8")


# ─── Image Generation (from PDF pages) ────────────────────────────────


def generate_page_image(
    pdf_path: Path, page_num: int, output_path: Path, dpi: int = 150
):
    """Export a full page as a PNG image."""
    png_bytes = extract_page_image(pdf_path, page_num, dpi)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_bytes(png_bytes)


def generate_region_image(
    pdf_path: Path,
    page_num: int,
    rect: tuple[float, float, float, float],
    output_path: Path,
    dpi: int = 200,
):
    """Export a page region as a PNG (simulates a diagram/drawing export)."""
    png_bytes = extract_page_region_image(pdf_path, page_num, rect, dpi)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_bytes(png_bytes)


def generate_tiff_scan(pdf_path: Path, page_num: int, output_path: Path, dpi: int = 200):
    """Render a page as TIFF (simulating a scanned legacy drawing)."""
    from PIL import Image
    import io

    png_bytes = extract_page_image(pdf_path, page_num, dpi)
    img = Image.open(io.BytesIO(png_bytes))
    # Convert to grayscale to simulate a scan
    img = img.convert("L")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(str(output_path), format="TIFF", compression="tiff_lzw")


# ═══════════════════════════════════════════════════════════════════════
# CORPUS GENERATION: HSI_SYS_015G
# ═══════════════════════════════════════════════════════════════════════


def generate_hsi_015g_corpus():
    """Generate all source documents for HSI_SYS_015G (Spectrometer ICD)."""
    pdf = ICDS_DIR / "HSI_SYS_015G.pdf"
    out = CORPUS_DIR / "hsi_sys_015g"

    # ─── DOCX ──────────────────────────────────────────────────────

    # Mechanical requirements (from page 5)
    mech_text = extract_page_text(pdf, 5)
    generate_docx(
        title="HESSI Spectrometer Mechanical Interface Requirements",
        sections=[
            ("Interface Drawing", "The mechanical configuration of the Spectrometer is shown in the Spectrometer ICD Drawing (reference 1). This document defines the mounting interface, envelope constraints, and alignment requirements."),
            ("Mass Properties", "Reference (6) shows the instrument mass properties, including current best estimate and maximum (with margin). The Spectrometer assembly mass shall not exceed the allocation defined in the system mass budget."),
            ("Field of View and Alignment", "The nine detectors must be aligned to be within 1mm of concentric to the field of view of each of the nine grid pairs on the imager. This alignment is achieved when installing the Spectrometer on the spacecraft by shimming at the imager interface points.\n\nAn all-sky field of view is also desired for the detectors. This means that reasonable effort should be made to minimize the amount of metal in this region. The center of this field of view is defined by an XY plane through the middle of the detectors, at spacecraft station Z=+7.32 inches. The field of view then extends radially outward 15 degrees above and below this plane, starting at the spacecraft Z axis."),
            ("Mechanisms", "2.4.1 Cryocooler\nThe Cryocooler contains a free piston, linear motion, Helium gas compressor that is powered by the CPC and is actively controlled by the IDPU. The cryocooler contains a counterbalance mass that is also controlled by the IDPU. The residual operating forces will not exceed 0.5 newtons driven at 59 Hz.\n\n2.4.2 Hi-Z Shutters\nThe Spectrometer includes two Hi-Z shutters used to reduce the science data rate by blocking out low energy radiation during periods of high flux."),
        ],
        output_path=out / "docx" / "HSI_Mech_Requirements_v2.docx",
        author="UCB Mechanical Engineering",
        doc_number="HSI-MECH-REQ-002",
        revision="v2",
        date="2024-09-10",
    )

    # Thermal test report (from page 6)
    thermal_text = extract_pages_text(pdf, range(6, 8))
    generate_docx(
        title="HESSI Spectrometer Thermal Vacuum Test Report",
        sections=[
            ("Test Objective", "Verify that the Spectrometer thermal design meets the temperature requirements specified in the Spacecraft to Spectrometer ICD (HSI_SYS_015G) Section 3.3. Demonstrate thermal margin under worst-case hot and cold orbital conditions."),
            ("Test Configuration", "The Spectrometer flight unit was installed in Thermal Vacuum Chamber 4 at SSL Berkeley. The test article was instrumented with 24 thermocouples at key interface locations. The chamber shroud was configured to simulate orbital thermal environments."),
            ("Results Summary", "All temperature limits were met with positive margin. The Cryostat achieved operating temperature within 48 hours of cooldown initiation. The bus side interface temperature remained within -10C to +40C throughout all test phases. The spectrometer side interface remained within -20C to +30C."),
            ("Heater Performance", "The programmable spectrometer heaters (up to 15W maximum) maintained the detectors above survival temperature during eclipse periods. Thermostat characteristics matched the specification in Table 3.2.1-1 of the ICD."),
        ],
        output_path=out / "docx" / "Thermal_Test_Report.docx",
        author="UCB Thermal Engineering",
        doc_number="HSI-TVT-RPT-001",
        revision="v1",
        date="2024-11-15",
    )

    # Integration procedure
    generate_docx(
        title="HESSI Spectrometer Integration Procedure",
        sections=[
            ("Scope", "This procedure defines the step-by-step process for integrating the Spectrometer assembly onto the HESSI spacecraft bus. It covers mechanical installation, electrical connections, alignment verification, and functional checkout."),
            ("Prerequisites", "- Spacecraft bus in integration configuration\n- Spectrometer assembly passed standalone functional test\n- Clean room environment (Class 10,000 or better)\n- Alignment GSE available and calibrated\n- Shimming materials available (0.001 to 0.020 inch range)"),
            ("Mechanical Installation", "Step 1: Position the Spectrometer above the spacecraft interface using the handling fixture.\nStep 2: Lower the Spectrometer onto the four mounting pad interface points.\nStep 3: Install the 8 mounting bolts finger-tight.\nStep 4: Perform preliminary alignment check using theodolite.\nStep 5: Install shims as required to achieve 1mm concentricity.\nStep 6: Torque all mounting bolts to 45 in-lbs."),
            ("Electrical Connections", "Step 7: Connect the power harness (J1 connector) to the CPC interface.\nStep 8: Connect the data harness (J2 connector) to the IDPU interface.\nStep 9: Connect the heater harness (J3 connector) to the thermal control bus.\nStep 10: Verify all connections with continuity check."),
            ("Alignment Verification", "Step 11: Perform final alignment survey using laser tracker.\nStep 12: Verify detector array concentricity is within 1mm of grid pair centers.\nStep 13: Record alignment data in traveler document."),
        ],
        output_path=out / "docx" / "Integration_Procedure_v1.docx",
        author="Spectrum Astro Integration Team",
        doc_number="HSI-INT-PROC-001",
        revision="v1",
        date="2025-03-01",
    )

    # ─── XLSX ──────────────────────────────────────────────────────

    # Thermostat parameters (Table 3.2.1-1)
    generate_xlsx(
        title="Spectrometer Thermostat Characteristics",
        headers=["Characteristic", "Value", "Units", "Tolerance", "Notes"],
        rows=[
            ["Thermostat Type", "Programmable", "—", "—", "Controlled by IDPU"],
            ["Maximum Power", "15", "W", "+0/-1W", "Per heater circuit"],
            ["Set Point Range", "-20 to +40", "°C", "±2°C", "Software configurable"],
            ["Hysteresis", "2", "°C", "±0.5°C", "Fixed by hardware"],
            ["Response Time", "<30", "sec", "—", "To reach set point"],
            ["Number of Channels", "4", "—", "—", "Independent thermostats"],
            ["Sensor Type", "Platinum RTD", "—", "—", "100 ohm at 0°C"],
            ["Control Loop", "PID", "—", "—", "Gains set during I&T"],
        ],
        output_path=out / "xlsx" / "Thermostat_Parameters.xlsx",
        sheet_name="Thermostat Chars",
        doc_ref="HSI_SYS_015G Table 3.2.1-1",
        date="2024-05-30",
    )

    # Mass properties
    generate_xlsx(
        title="HSI Spectrometer Mass Properties",
        headers=["Component", "CBE Mass (kg)", "Margin (%)", "Max Mass (kg)", "CG X (in)", "CG Y (in)", "CG Z (in)"],
        rows=[
            ["Cryostat Assembly", "28.5", "10", "31.4", "0.0", "0.0", "+7.32"],
            ["Detector Array (9x)", "4.2", "5", "4.4", "0.0", "0.0", "+7.50"],
            ["Cryocooler", "6.8", "15", "7.8", "+2.1", "0.0", "+5.00"],
            ["Hi-Z Shutters (2x)", "1.06", "10", "1.17", "0.0", "0.0", "+11.30"],
            ["Thermal Hardware", "2.1", "20", "2.5", "0.0", "0.0", "+6.80"],
            ["Harnesses", "1.8", "15", "2.1", "—", "—", "—"],
            ["TOTAL", "44.5", "11.3", "49.4", "0.2", "0.0", "+7.15"],
        ],
        output_path=out / "xlsx" / "HSI_Mass_Properties.xlsx",
        sheet_name="Mass Budget",
        doc_ref="HSI_SYS_015G Section 2.2",
        date="2025-04-20",
    )

    # Power budget
    generate_xlsx(
        title="HSI Spectrometer Electrical Power Budget",
        headers=["Mode", "Subsystem", "Voltage (V)", "Current (A)", "Power (W)", "Duty Cycle (%)", "Avg Power (W)"],
        rows=[
            ["Science", "Detectors", "5.0", "0.8", "4.0", "100", "4.0"],
            ["Science", "Cryocooler", "28.0", "2.5", "70.0", "100", "70.0"],
            ["Science", "IDPU", "5.0", "2.0", "10.0", "100", "10.0"],
            ["Science", "Heaters", "28.0", "0.54", "15.0", "30", "4.5"],
            ["Eclipse", "Survival Heaters", "28.0", "0.54", "15.0", "100", "15.0"],
            ["Eclipse", "IDPU Standby", "5.0", "0.5", "2.5", "100", "2.5"],
            ["Safehold", "Essential Heaters", "28.0", "0.25", "7.0", "100", "7.0"],
        ],
        output_path=out / "xlsx" / "Power_Budget_v2.xlsx",
        sheet_name="Power Budget",
        doc_ref="HSI_SYS_015G Section 4",
        date="2024-10-05",
    )

    # ─── PPTX ──────────────────────────────────────────────────────

    generate_pptx(
        title="HESSI Spectrometer Thermal Analysis",
        subtitle="Critical Design Review",
        slides=[
            ("Thermal Design Overview", "Radiative and conductive heat transfer between Cryostat, spacecraft, and space\nDesign meets thermal constraints in Section 3.3\nPrimary dissipation via large bottom radiator surface\nConductive coupling through aluminum mounting flanges"),
            ("Thermal Requirements", "Bus Side Interface: -10°C to +40°C\nSpectrometer Side Interface: -20°C to +30°C\nDetector Operating: -198°C to -190°C\nCryocooler Reject: < +50°C"),
            ("Thermal Model Results", "Worst Case Hot: All interfaces within limits with 5°C margin\nWorst Case Cold: Heaters maintain survival with 15W allocation\nTransient Eclipse: Detector temp stable within 2°C"),
            ("Heater Design", "4 independent programmable thermostats\nMaximum 15W per channel\nPID control with flight-adjustable gains\nPlatinum RTD sensors (100 ohm at 0°C)"),
            ("Spectrometer Power Dissipation", "Science Mode: 88.5W total\nEclipse Mode: 17.5W (heaters + standby)\nDissipation primarily through bottom radiator\nConducted path: 2W maximum via mounting flanges"),
        ],
        output_path=out / "pptx" / "Thermal_Analysis_CDR.pptx",
        date="2024-02-10",
        author="UCB Thermal Engineering",
    )

    # ─── HTML (Confluence page) ────────────────────────────────────

    generate_html(
        title="HSI Power Interface Specification",
        sections=[
            ("Power Bus Interface", "The Spectrometer receives unregulated 28V power from the spacecraft bus through connector J1. The power interface provides over-current protection via a 5A fuse on the spacecraft side. The Spectrometer internal power conditioning unit (CPC) regulates 28V to the required secondary voltages (5V, ±12V, -190V detector bias)."),
            ("Power Modes", "Science Mode: 88.5W average (cryocooler dominant)\nEclipse Mode: 17.5W (survival heaters + IDPU standby)\nSafehold Mode: 7.0W (essential heaters only)\n\nThe spacecraft must provide continuous power in all modes. Power cycling the Spectrometer requires 48+ hours for cryocooler re-cooldown."),
            ("Current Requirements", "Inrush current at power-on: < 10A for < 100ms\nSteady-state science mode: 3.2A at 28V\nPeak transient (shutter actuation): 5A for < 50ms\n\nThe spacecraft bus must accommodate these transient loads without voltage droop below 24V."),
            ("Open Items", "TBR-UCB-102: Cryocooler residual operating forces (currently 0.5N, awaiting final vibration test)\nTBR-UCB-110: Shutter actuation duration (currently 0.5s, design may change to 0.3s)"),
        ],
        output_path=out / "html" / "HSI_Power_ICD_Wiki.html",
        space="HESSI Engineering",
        last_modified="2025-05-30",
    )

    # ─── Images ────────────────────────────────────────────────────

    # Spectrometer mounting region (page 5, top half — mechanical section)
    generate_region_image(
        pdf, 5, (50, 80, 560, 400), out / "images" / "Spectrometer_Mount_Drawing.png"
    )

    # Thermal table region (page 7 — thermostat table)
    generate_region_image(
        pdf, 7, (50, 200, 560, 450), out / "images" / "Thermal_Table_Extract.png"
    )

    # Full page as TIFF "scanned drawing" (page 5 — mechanical section with dense text)
    generate_tiff_scan(pdf, 5, out / "images" / "Connector_Pinout_J1.tiff")

    # Block diagram (page 4 — introduction with system overview)
    generate_region_image(
        pdf, 4, (50, 300, 560, 700), out / "images" / "System_Overview_Diagram.png"
    )

    print(f"  Generated HSI_SYS_015G corpus in {out}")


# ═══════════════════════════════════════════════════════════════════════
# CORPUS GENERATION: TSAFE (20130010957)
# ═══════════════════════════════════════════════════════════════════════


def generate_tsafe_corpus():
    """Generate source documents for TSAFE ICD."""
    pdf = ICDS_DIR / "20130010957.pdf"
    out = CORPUS_DIR / "tsafe"

    # Data format spec
    data_text = extract_pages_text(pdf, range(7, 12))
    generate_docx(
        title="TSAFE Data Format Specification",
        sections=[
            ("Input Data Formats", "The input data format from radar consists of a line of standard ASCII text with data fields delimited by spaces. Multiple records can be sent in a single message by separating the records with semicolons.\n\nThree record types trigger a conflict check: Track Update (TRK), Vector Amendment (VEC), and Altitude Amendment (ALT)."),
            ("Track Update Record (TRK)", "Each Track Update record triggers an immediate check for conflicts between the flight for which the record applies and all other traffic in the Center. Fields include: record type, aircraft ID, position (latitude, longitude, altitude), ground speed, heading, vertical rate, and timestamp."),
            ("Output Message Format", "TSAFE advisory messages are formatted as ASCII text with field delimiters. Each advisory includes: conflict pair identification, time to loss of separation, closest point of approach, and recommended resolution maneuver."),
        ],
        output_path=out / "docx" / "TSAFE_Message_Spec.docx",
        author="NASA Ames Research Center",
        doc_number="TSAFE-DFS-002",
        revision="v2",
        date="2012-09-15",
    )

    # Test results
    generate_xlsx(
        title="TSAFE Performance Test Matrix",
        headers=["Test Case", "Scenario", "Expected (s)", "Actual (s)", "Margin (%)", "Status"],
        rows=[
            ["TC-001", "Single conflict pair", "30", "22", "26.7", "PASS"],
            ["TC-002", "Multiple simultaneous conflicts", "30", "28", "6.7", "PASS"],
            ["TC-003", "False alarm rate target", "5.0", "3.2", "36.0", "PASS"],
            ["TC-004", "Missed detection rate", "1.0", "0.4", "60.0", "PASS"],
            ["TC-005", "Latency under load (100 aircraft)", "2.0", "1.4", "30.0", "PASS"],
            ["TC-006", "Recovery after 30s data gap", "10", "8", "20.0", "PASS"],
        ],
        output_path=out / "xlsx" / "TSAFE_Test_Results.xlsx",
        sheet_name="Test Matrix",
        doc_ref="20130010957 Section 5",
        date="2013-04-15",
    )

    # Architecture overview presentation
    generate_pptx(
        title="TSAFE System Architecture",
        subtitle="Traffic Safety Advisor for Controllers",
        slides=[
            ("System Overview", "TSAFE provides short-term conflict advisories to air traffic controllers\nOperates on live radar track data\nPredicts trajectories 3-5 minutes ahead\nDetects loss-of-separation violations"),
            ("Input/Output Architecture", "Inputs: Radar track updates, flight plan amendments, weather data\nProcessing: Trajectory prediction + conflict detection algorithms\nOutputs: Advisory messages to controller displays\nInterface: Standard ASCII message format over TCP/IP"),
            ("Performance Characteristics", "Detection rate: > 99% of conflicts with > 60s warning\nFalse alarm rate: < 5% of total advisories\nProcessing latency: < 2 seconds for 100 aircraft\nUpdate rate: Every radar scan (4.8 seconds)"),
        ],
        output_path=out / "pptx" / "TSAFE_Architecture_Overview.pptx",
        date="2012-03-10",
        author="NASA Ames TSAFE Team",
    )

    # System diagram (page 7 — has data format content)
    generate_region_image(
        pdf, 7, (50, 80, 560, 500), out / "images" / "TSAFE_Architecture_Diagram.png"
    )

    # Data flow (page 7)
    generate_region_image(
        pdf, 7, (50, 80, 560, 400), out / "images" / "Radar_Interface_Diagram.png"
    )

    print(f"  Generated TSAFE corpus in {out}")


# ═══════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════


def main():
    print("Generating test corpus from ICD documents...")
    print()

    if not ICDS_DIR.exists():
        print(f"ERROR: ICD directory not found: {ICDS_DIR}")
        return

    generate_hsi_015g_corpus()
    generate_tsafe_corpus()

    # Manifest (basic — will be extended)
    manifest = {
        "generated_at": datetime.now().isoformat(),
        "source_icds": [p.name for p in sorted(ICDS_DIR.glob("*.pdf")) if p.stat().st_size > 200],
        "corpus_root": str(CORPUS_DIR),
        "note": "Generated by test_corpus/scripts/generate_corpus.py from real ICD content",
    }
    MANIFEST_PATH.write_text(json.dumps(manifest, indent=2), encoding="utf-8")
    print(f"\n  Manifest written to {MANIFEST_PATH}")
    print("\nDone.")


if __name__ == "__main__":
    main()
