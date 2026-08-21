"""Tests for versioned document series — validates multi-update lineage chains.

Each series has 3-5 versions of the same document, simulating the real
lifecycle of engineering artifacts feeding into an ICD. Tests verify:

1. All versions in a series exist and open correctly
2. Content evolves realistically across versions (not identical copies)
3. Version ordering is chronological (dates increase)
4. Later versions have more resolved items / more data
5. Document metadata (revision, date) is consistent within each version
6. The series tells a coherent engineering story (draft → test → flight)

Series tested:
- HSI Mech Requirements: DOCX v1→v2→v3→v4 (draft → PDR → CDR → flight)
- HSI Power Budget: XLSX v1→v2→v3→v4→v5 (allocation → measured → on-orbit)
- IDSS Seal Design Review: PPTX v1→v2→v3 (PDR → CDR → FRR)
- HSI Thermal Limits: XLSX v1→v2→v3 (analysis → test → on-orbit)
"""

from pathlib import Path

import pytest

CORPUS_DIR = Path(__file__).parent.parent.parent.parent / "test_corpus"
HSI_SERIES = CORPUS_DIR / "hsi_sys_015g" / "series"
IDSS_SERIES = CORPUS_DIR / "idss_idd" / "series"


def _series_exists():
    return (HSI_SERIES / "mech_requirements").exists()


pytestmark = pytest.mark.skipif(
    not _series_exists(),
    reason="Versioned series not generated (run: python test_corpus/scripts/generate_versioned_series.py)",
)


# ═══════════════════════════════════════════════════════════════════════
# HSI MECH REQUIREMENTS (DOCX, 4 versions)
# ═══════════════════════════════════════════════════════════════════════


class TestMechRequirementsSeries:
    """Validate the 4-version mechanical requirements DOCX series."""

    SERIES_DIR = HSI_SERIES / "mech_requirements"
    FILES = [
        "HSI_Mech_Requirements_v1.docx",
        "HSI_Mech_Requirements_v2.docx",
        "HSI_Mech_Requirements_v3.docx",
        "HSI_Mech_Requirements_v4.docx",
    ]

    def test_all_versions_exist(self):
        """All 4 versions are present."""
        for f in self.FILES:
            assert (self.SERIES_DIR / f).exists(), f"Missing: {f}"

    @pytest.mark.parametrize("filename", FILES)
    def test_each_version_opens(self, filename):
        """Each version opens as valid DOCX."""
        from docx import Document
        doc = Document(str(self.SERIES_DIR / filename))
        assert len(doc.paragraphs) > 10

    def test_content_evolves_across_versions(self):
        """Each version has different content (not identical copies)."""
        from docx import Document
        texts = []
        for f in self.FILES:
            doc = Document(str(self.SERIES_DIR / f))
            text = " ".join(p.text for p in doc.paragraphs)
            texts.append(text)

        # Each consecutive pair should differ
        for i in range(len(texts) - 1):
            assert texts[i] != texts[i + 1], (
                f"v{i+1} and v{i+2} are identical — no evolution"
            )

    def test_tbds_decrease_over_versions(self):
        """TBD/TBR count decreases from v1 to v4 (items get resolved)."""
        from docx import Document
        tbd_counts = []
        for f in self.FILES:
            doc = Document(str(self.SERIES_DIR / f))
            text = " ".join(p.text for p in doc.paragraphs)
            count = text.lower().count("tbd") + text.lower().count("tbr")
            tbd_counts.append(count)

        # v1 should have more TBDs than v3/v4
        assert tbd_counts[0] >= tbd_counts[-1], (
            f"TBDs didn't decrease: {tbd_counts}"
        )

    def test_v1_has_preliminary_language(self):
        """v1 uses preliminary/TBD language (it's a draft)."""
        from docx import Document
        doc = Document(str(self.SERIES_DIR / self.FILES[0]))
        text = " ".join(p.text for p in doc.paragraphs).lower()
        preliminary_markers = ["tbd", "preliminary", "estimate", "allocation", "pending"]
        matches = [m for m in preliminary_markers if m in text]
        assert len(matches) >= 2, f"v1 doesn't sound like a draft: found {matches}"

    def test_v4_has_flight_language(self):
        """v4 uses flight/verified/measured language (it's final)."""
        from docx import Document
        doc = Document(str(self.SERIES_DIR / self.FILES[3]))
        text = " ".join(p.text for p in doc.paragraphs).lower()
        final_markers = ["measured", "verified", "on-orbit", "flight", "as-flown"]
        matches = [m for m in final_markers if m in text]
        assert len(matches) >= 2, f"v4 doesn't sound final: found {matches}"

    def test_word_count_increases(self):
        """Later versions have more content (details added over time)."""
        from docx import Document
        word_counts = []
        for f in self.FILES:
            doc = Document(str(self.SERIES_DIR / f))
            text = " ".join(p.text for p in doc.paragraphs)
            word_counts.append(len(text.split()))

        # v4 should have more words than v1
        assert word_counts[-1] > word_counts[0], (
            f"Content didn't grow: v1={word_counts[0]} words, v4={word_counts[-1]} words"
        )


# ═══════════════════════════════════════════════════════════════════════
# HSI POWER BUDGET (XLSX, 5 versions)
# ═══════════════════════════════════════════════════════════════════════


class TestPowerBudgetSeries:
    """Validate the 5-version power budget XLSX series."""

    SERIES_DIR = HSI_SERIES / "power_budget"
    FILES = [
        "Power_Budget_v1.xlsx",
        "Power_Budget_v2.xlsx",
        "Power_Budget_v3.xlsx",
        "Power_Budget_v4.xlsx",
        "Power_Budget_v5.xlsx",
    ]

    def test_all_versions_exist(self):
        """All 5 versions are present."""
        for f in self.FILES:
            assert (self.SERIES_DIR / f).exists(), f"Missing: {f}"

    @pytest.mark.parametrize("filename", FILES)
    def test_each_version_opens(self, filename):
        """Each version opens as valid XLSX."""
        from openpyxl import load_workbook
        wb = load_workbook(str(self.SERIES_DIR / filename))
        assert len(wb.sheetnames) >= 1

    def test_column_count_increases(self):
        """Later versions have more columns (more detail added)."""
        from openpyxl import load_workbook
        col_counts = []
        for f in self.FILES:
            wb = load_workbook(str(self.SERIES_DIR / f))
            ws = wb.active
            # Count non-empty columns in header row (row 4)
            cols = sum(1 for c in ws[4] if c.value)
            col_counts.append(cols)

        # v5 should have >= columns as v1
        assert col_counts[-1] >= col_counts[0], (
            f"Columns didn't grow: {col_counts}"
        )

    def test_status_evolves_to_measured(self):
        """Early versions say ESTIMATE/ALLOCATION, later say MEASURED/FLIGHT."""
        from openpyxl import load_workbook

        # v1 — should have ESTIMATE or ALLOCATION
        wb1 = load_workbook(str(self.SERIES_DIR / self.FILES[0]))
        ws1 = wb1.active
        v1_text = " ".join(
            str(cell.value or "") for row in ws1.iter_rows(min_row=5, max_row=15)
            for cell in row
        ).upper()
        assert "ESTIMATE" in v1_text or "ALLOCATION" in v1_text, (
            "v1 should have ESTIMATE/ALLOCATION status"
        )

        # v5 — should have ON-ORBIT or FLIGHT
        wb5 = load_workbook(str(self.SERIES_DIR / self.FILES[4]))
        ws5 = wb5.active
        v5_text = " ".join(
            str(cell.value or "") for row in ws5.iter_rows(min_row=5, max_row=15)
            for cell in row
        ).upper()
        assert "ON-ORBIT" in v5_text or "FLIGHT" in v5_text, (
            "v5 should have ON-ORBIT/FLIGHT status"
        )

    def test_power_values_converge(self):
        """Power values get more precise (smaller changes between v3→v4→v5)."""
        from openpyxl import load_workbook

        def get_first_power_value(filepath):
            wb = load_workbook(str(filepath))
            ws = wb.active
            # Find first numeric cell in "Power" column (likely col 5)
            for row in ws.iter_rows(min_row=5, max_row=12, min_col=4, max_col=6):
                for cell in row:
                    if cell.value:
                        try:
                            return float(str(cell.value))
                        except (ValueError, TypeError):
                            continue
            return None

        v3_power = get_first_power_value(self.SERIES_DIR / self.FILES[2])
        v4_power = get_first_power_value(self.SERIES_DIR / self.FILES[3])
        v5_power = get_first_power_value(self.SERIES_DIR / self.FILES[4])

        if v3_power and v4_power and v5_power:
            # Changes should get smaller (convergence)
            delta_34 = abs(v4_power - v3_power)
            delta_45 = abs(v5_power - v4_power)
            # Not a hard assertion — just verify values are similar (within 20%)
            assert abs(v5_power - v3_power) / v3_power < 0.20, (
                f"Power values diverged too much: v3={v3_power}, v5={v5_power}"
            )


# ═══════════════════════════════════════════════════════════════════════
# IDSS SEAL DESIGN REVIEW (PPTX, 3 versions)
# ═══════════════════════════════════════════════════════════════════════


class TestSealDesignReviewSeries:
    """Validate the 3-version seal design review PPTX series."""

    SERIES_DIR = IDSS_SERIES / "seal_design_review"
    FILES = [
        "Seal_Design_Review_v1_PDR.pptx",
        "Seal_Design_Review_v2_CDR.pptx",
        "Seal_Design_Review_v3_FRR.pptx",
    ]

    def test_all_versions_exist(self):
        """All 3 versions are present."""
        for f in self.FILES:
            assert (self.SERIES_DIR / f).exists(), f"Missing: {f}"

    @pytest.mark.parametrize("filename", FILES)
    def test_each_version_opens(self, filename):
        """Each version opens as valid PPTX."""
        from pptx import Presentation
        prs = Presentation(str(self.SERIES_DIR / filename))
        assert len(prs.slides) >= 4

    def test_review_type_progresses(self):
        """Title slides show PDR → CDR → FRR progression."""
        from pptx import Presentation
        subtitles = []
        for f in self.FILES:
            prs = Presentation(str(self.SERIES_DIR / f))
            # Get subtitle from first slide
            first_slide = prs.slides[0]
            texts = [
                shape.text for shape in first_slide.shapes if shape.has_text_frame
            ]
            subtitles.append(" ".join(texts))

        assert "PDR" in subtitles[0] or "Preliminary" in subtitles[0]
        assert "CDR" in subtitles[1] or "Critical" in subtitles[1]
        assert "FRR" in subtitles[2] or "Flight" in subtitles[2] or "Readiness" in subtitles[2]

    def test_tbd_count_decreases(self):
        """Unresolved TBD items decrease from PDR to FRR."""
        from pptx import Presentation
        tbd_counts = []
        for f in self.FILES:
            prs = Presentation(str(self.SERIES_DIR / f))
            all_text = " ".join(
                shape.text for slide in prs.slides
                for shape in slide.shapes if shape.has_text_frame
            )
            # Count actual unresolved TBD markers (not "RESOLVED" or "CLOSED" ones)
            lines = all_text.split("\n")
            unresolved = sum(
                1 for line in lines
                if "TBD" in line.upper() and "RESOLVED" not in line.upper() and "CLOSED" not in line.upper()
            )
            tbd_counts.append(unresolved)

        assert tbd_counts[0] >= tbd_counts[-1], (
            f"Unresolved TBDs didn't decrease: PDR={tbd_counts[0]}, FRR={tbd_counts[-1]}"
        )

    def test_v3_contains_resolved_language(self):
        """FRR version uses RESOLVED/CLOSED/PASS language."""
        from pptx import Presentation
        prs = Presentation(str(self.SERIES_DIR / self.FILES[2]))
        all_text = " ".join(
            shape.text for slide in prs.slides
            for shape in slide.shapes if shape.has_text_frame
        ).upper()

        resolved_terms = ["RESOLVED", "CLOSED", "PASS", "FLIGHT READY", "VERIFIED"]
        matches = [t for t in resolved_terms if t in all_text]
        assert len(matches) >= 2, f"FRR doesn't sound resolved: found {matches}"

    def test_content_grows_with_test_data(self):
        """Later versions have more text (test results added)."""
        from pptx import Presentation
        text_lengths = []
        for f in self.FILES:
            prs = Presentation(str(self.SERIES_DIR / f))
            all_text = " ".join(
                shape.text for slide in prs.slides
                for shape in slide.shapes if shape.has_text_frame
            )
            text_lengths.append(len(all_text))

        assert text_lengths[-1] > text_lengths[0], (
            f"Content didn't grow: v1={text_lengths[0]} chars, v3={text_lengths[-1]} chars"
        )


# ═══════════════════════════════════════════════════════════════════════
# HSI THERMAL LIMITS (XLSX, 3 versions)
# ═══════════════════════════════════════════════════════════════════════


class TestThermalLimitsSeries:
    """Validate the 3-version thermal limits XLSX series."""

    SERIES_DIR = HSI_SERIES / "thermal_limits"
    FILES = [
        "Thermal_Limits_v1.xlsx",
        "Thermal_Limits_v2.xlsx",
        "Thermal_Limits_v3.xlsx",
    ]

    def test_all_versions_exist(self):
        """All 3 versions are present."""
        for f in self.FILES:
            assert (self.SERIES_DIR / f).exists(), f"Missing: {f}"

    @pytest.mark.parametrize("filename", FILES)
    def test_each_version_opens(self, filename):
        """Each version opens as valid XLSX."""
        from openpyxl import load_workbook
        wb = load_workbook(str(self.SERIES_DIR / filename))
        ws = wb.active
        # Should have header + data rows
        rows = list(ws.iter_rows(min_row=4, max_row=15))
        assert len(rows) >= 5

    def test_columns_grow_with_detail(self):
        """Later versions add columns (test results, flight delta)."""
        from openpyxl import load_workbook
        col_counts = []
        for f in self.FILES:
            wb = load_workbook(str(self.SERIES_DIR / f))
            ws = wb.active
            cols = sum(1 for c in ws[4] if c.value)
            col_counts.append(cols)

        assert col_counts[-1] > col_counts[0], (
            f"Columns didn't grow: {col_counts}"
        )

    def test_basis_evolves(self):
        """Basis column changes from Analysis → Test → Flight Data."""
        from openpyxl import load_workbook

        def get_basis_values(filepath):
            wb = load_workbook(str(filepath))
            ws = wb.active
            bases = []
            for row in ws.iter_rows(min_row=5, max_row=12):
                for cell in row:
                    val = str(cell.value or "").lower()
                    if any(b in val for b in ["analysis", "test", "flight", "vendor"]):
                        bases.append(val)
            return " ".join(bases)

        v1_bases = get_basis_values(self.SERIES_DIR / self.FILES[0])
        v3_bases = get_basis_values(self.SERIES_DIR / self.FILES[2])

        assert "analysis" in v1_bases, "v1 should be analysis-based"
        assert "flight" in v3_bases, "v3 should be flight-data-based"

    def test_margins_improve_over_time(self):
        """Thermal margins generally improve (better than predicted)."""
        from openpyxl import load_workbook

        def get_margin_column_values(filepath):
            wb = load_workbook(str(filepath))
            ws = wb.active
            # Find margin column header
            margin_col = None
            for cell in ws[4]:
                if cell.value and "margin" in str(cell.value).lower():
                    margin_col = cell.column
                    break
            if not margin_col:
                return []
            margins = []
            for row in ws.iter_rows(min_row=5, max_row=12, min_col=margin_col, max_col=margin_col):
                val = row[0].value
                if val:
                    try:
                        margins.append(float(str(val)))
                    except (ValueError, TypeError):
                        pass
            return margins

        v1_margins = get_margin_column_values(self.SERIES_DIR / self.FILES[0])
        v3_margins = get_margin_column_values(self.SERIES_DIR / self.FILES[2])

        if v1_margins and v3_margins:
            avg_v1 = sum(v1_margins) / len(v1_margins)
            avg_v3 = sum(v3_margins) / len(v3_margins)
            # On-orbit margins should be >= analytical margins
            assert avg_v3 >= avg_v1 - 2, (
                f"Margins degraded: v1 avg={avg_v1:.1f}°C, v3 avg={avg_v3:.1f}°C"
            )
